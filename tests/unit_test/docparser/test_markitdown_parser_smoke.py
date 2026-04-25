import json
import zipfile
from collections.abc import Callable
from pathlib import Path

import pytest
from openpyxl import Workbook
from pptx import Presentation

from aperag.docparser.doc_parser import DocParser


def _write_text(path: Path, content: str) -> None:
    path.write_text(content, encoding="utf-8")


def _write_long_markdown(path: Path) -> None:
    sections = [
        "# Long Markdown Smoke",
        "",
        "This synthetic document exercises the default parser path with a longer markdown body.",
        "",
    ]
    for idx in range(1, 7):
        sections.extend(
            [
                f"## Section {idx}",
                f"Section {idx} body line for parser regression coverage.",
                f"- bullet {idx}a",
                f"- bullet {idx}b",
                "",
            ]
        )
    sections.extend(
        [
            "| Region | Revenue | Notes |",
            "| --- | ---: | --- |",
            "| East | 120 | markdown parser table |",
            "| West | 140 | long markdown parser tail |",
        ]
    )
    path.write_text("\n".join(sections), encoding="utf-8")


def _write_long_html(path: Path) -> None:
    sections = []
    for idx, region in enumerate(["Alpha", "Beta", "Gamma", "Delta"], start=1):
        sections.append(f"<section><h2>HTML Section {idx}</h2><p>{region} html parser section body</p></section>")
    content = """
<html>
  <body>
    <h1>Long HTML Smoke</h1>
    <p>html parser table note</p>
    <table>
      <tr><th>Region</th><th>Status</th></tr>
      <tr><td>Region Alpha</td><td>ready</td></tr>
      <tr><td>Region Delta</td><td>watch</td></tr>
    </table>
    {sections}
  </body>
</html>
""".strip().format(sections="".join(sections))
    path.write_text(content, encoding="utf-8")


def _write_ipynb(path: Path) -> None:
    notebook = {
        "cells": [
            {
                "cell_type": "markdown",
                "metadata": {},
                "source": ["# Notebook Smoke\n", "ipynb parser smoke"],
            },
            {
                "cell_type": "code",
                "execution_count": 1,
                "metadata": {},
                "outputs": [],
                "source": ['print("hello notebook")'],
            },
        ],
        "metadata": {},
        "nbformat": 4,
        "nbformat_minor": 5,
    }
    path.write_text(json.dumps(notebook), encoding="utf-8")


def _write_docx(path: Path) -> None:
    with zipfile.ZipFile(path, "w") as archive:
        archive.writestr(
            "[Content_Types].xml",
            """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>
</Types>""",
        )
        archive.writestr(
            "_rels/.rels",
            """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>
</Relationships>""",
        )
        archive.writestr(
            "word/document.xml",
            """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<w:document xmlns:wpc="http://schemas.microsoft.com/office/word/2010/wordprocessingCanvas" xmlns:mc="http://schemas.openxmlformats.org/markup-compatibility/2006" xmlns:o="urn:schemas-microsoft-com:office:office" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:m="http://schemas.openxmlformats.org/officeDocument/2006/math" xmlns:v="urn:schemas-microsoft-com:vml" xmlns:wp14="http://schemas.microsoft.com/office/word/2010/wordprocessingDrawing" xmlns:wp="http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing" xmlns:w10="urn:schemas-microsoft-com:office:word" xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main" xmlns:w14="http://schemas.microsoft.com/office/word/2010/wordml" xmlns:wpg="http://schemas.microsoft.com/office/word/2010/wordprocessingGroup" xmlns:wpi="http://schemas.microsoft.com/office/word/2010/wordprocessingInk" xmlns:wne="http://schemas.microsoft.com/office/word/2006/wordml" xmlns:wps="http://schemas.microsoft.com/office/word/2010/wordprocessingShape" mc:Ignorable="w14 wp14">
  <w:body>
    <w:p><w:r><w:t>Docx Smoke Title</w:t></w:r></w:p>
    <w:p><w:r><w:t>Docx parser smoke body</w:t></w:r></w:p>
    <w:sectPr><w:pgSz w:w="12240" w:h="15840"/><w:pgMar w:top="1440" w:right="1440" w:bottom="1440" w:left="1440" w:header="708" w:footer="708" w:gutter="0"/></w:sectPr>
  </w:body>
</w:document>""",
        )


def _write_xlsx(path: Path) -> None:
    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = "Smoke"
    worksheet["A1"] = "Revenue"
    worksheet["B1"] = 42
    worksheet["A2"] = "Notes"
    worksheet["B2"] = "xlsx parser smoke"
    workbook.save(path)


def _write_wide_xlsx(path: Path) -> None:
    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = "Summary"
    headers = ["Region", "Q1", "Q2", "Q3", "Q4", "Owner", "Notes"]
    worksheet.append(headers)
    for idx in range(1, 16):
        worksheet.append(
            [
                f"Region-{idx}",
                idx * 10,
                idx * 12,
                idx * 14,
                idx * 16,
                f"owner-{idx}",
                f"row-{idx}-note",
            ]
        )

    details = workbook.create_sheet("Details")
    details.append(["Metric", "Value"])
    details.append(["North", "stable"])
    details.append(["South", "watch"])
    workbook.save(path)


def _write_pptx(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "PPTX Smoke Title"
    slide.placeholders[1].text = "pptx parser smoke body"
    presentation.save(path)


def _write_multi_slide_pptx(path: Path) -> None:
    presentation = Presentation()

    title_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    title_slide.shapes.title.text = "PPTX Deck Smoke"
    title_slide.placeholders[1].text = "Agenda slide parser smoke"

    bullet_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    bullet_slide.shapes.title.text = "Execution Plan"
    bullet_slide.placeholders[1].text = "Milestone A\nMilestone B\nMilestone C"

    closing_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    closing_slide.shapes.title.text = "Closing Summary"
    closing_slide.placeholders[1].text = "Closing slide parser smoke"
    presentation.save(path)


def _write_pdf(path: Path) -> None:
    stream = b"BT /F1 18 Tf 72 720 Td (PDF Parser Smoke) Tj ET"
    objects = [
        b"1 0 obj\n<< /Type /Catalog /Pages 2 0 R >>\nendobj\n",
        b"2 0 obj\n<< /Type /Pages /Kids [3 0 R] /Count 1 >>\nendobj\n",
        b"3 0 obj\n<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>\nendobj\n",
        b"4 0 obj\n<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>\nendobj\n",
        f"5 0 obj\n<< /Length {len(stream)} >>\nstream\n".encode("latin-1") + stream + b"\nendstream\nendobj\n",
    ]

    pdf = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for obj in objects:
        offsets.append(len(pdf))
        pdf.extend(obj)

    startxref = len(pdf)
    pdf.extend(f"xref\n0 {len(offsets)}\n".encode("latin-1"))
    pdf.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        pdf.extend(f"{offset:010d} 00000 n \n".encode("latin-1"))
    pdf.extend(f"trailer\n<< /Size {len(offsets)} /Root 1 0 R >>\nstartxref\n{startxref}\n%%EOF\n".encode("latin-1"))
    path.write_bytes(pdf)


def _write_complex_pdf(path: Path) -> None:
    stream1 = b"BT /F1 18 Tf 72 720 Td (Complex PDF Smoke) Tj 0 -24 Td (First page parser smoke body) Tj ET"
    stream2 = b"BT /F1 18 Tf 72 720 Td (Second Page Title) Tj 0 -24 Td (Second page parser smoke body) Tj ET"
    objects = [
        b"1 0 obj\n<< /Type /Catalog /Pages 2 0 R >>\nendobj\n",
        b"2 0 obj\n<< /Type /Pages /Kids [3 0 R 4 0 R] /Count 2 >>\nendobj\n",
        b"3 0 obj\n<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Resources << /Font << /F1 5 0 R >> >> /Contents 6 0 R >>\nendobj\n",
        b"4 0 obj\n<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Resources << /Font << /F1 5 0 R >> >> /Contents 7 0 R >>\nendobj\n",
        b"5 0 obj\n<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>\nendobj\n",
        f"6 0 obj\n<< /Length {len(stream1)} >>\nstream\n".encode("latin-1") + stream1 + b"\nendstream\nendobj\n",
        f"7 0 obj\n<< /Length {len(stream2)} >>\nstream\n".encode("latin-1") + stream2 + b"\nendstream\nendobj\n",
    ]

    pdf = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for obj in objects:
        offsets.append(len(pdf))
        pdf.extend(obj)

    startxref = len(pdf)
    pdf.extend(f"xref\n0 {len(offsets)}\n".encode("latin-1"))
    pdf.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        pdf.extend(f"{offset:010d} 00000 n \n".encode("latin-1"))
    pdf.extend(f"trailer\n<< /Size {len(offsets)} /Root 1 0 R >>\nstartxref\n{startxref}\n%%EOF\n".encode("latin-1"))
    path.write_bytes(pdf)


def _extract_text(parts: list) -> str:
    return "\n".join(part.content for part in parts if getattr(part, "content", None))


@pytest.mark.parametrize(
    ("filename", "writer", "expected_snippets"),
    [
        (
            "smoke.txt",
            lambda path: _write_text(path, "TXT parser smoke body"),
            ["TXT parser smoke body"],
        ),
        (
            "smoke.md",
            lambda path: _write_text(path, "# MD Smoke\nmarkdown parser smoke body"),
            ["MD Smoke", "markdown parser smoke body"],
        ),
        (
            "smoke.html",
            lambda path: _write_text(path, "<h1>HTML Smoke</h1><p>html parser smoke body</p>"),
            ["HTML Smoke", "html parser smoke body"],
        ),
        (
            "smoke.ipynb",
            _write_ipynb,
            ["Notebook Smoke", "ipynb parser smoke", 'print("hello notebook")'],
        ),
        (
            "smoke.docx",
            _write_docx,
            ["Docx Smoke Title", "Docx parser smoke body"],
        ),
        (
            "smoke.xlsx",
            _write_xlsx,
            ["Smoke", "Revenue", "xlsx parser smoke"],
        ),
        (
            "smoke.pptx",
            _write_pptx,
            ["PPTX Smoke Title", "pptx parser smoke body"],
        ),
        (
            "smoke.pdf",
            _write_pdf,
            ["PDF Parser Smoke"],
        ),
    ],
)
def test_doc_parser_smoke_for_common_document_formats(
    tmp_path: Path,
    filename: str,
    writer: Callable[[Path], None],
    expected_snippets: list[str],
) -> None:
    sample = tmp_path / filename
    writer(sample)

    parser = DocParser(parser_config={"use_markitdown": True, "use_mineru": False})
    assert parser.accept(sample.suffix)

    parts = parser.parse_file(sample, metadata={"source": filename})
    combined_text = _extract_text(parts)

    assert parts, f"{filename} should produce parsed parts"
    for snippet in expected_snippets:
        assert snippet in combined_text, f"{filename} should contain {snippet!r}, got: {combined_text!r}"


@pytest.mark.parametrize(
    ("filename", "writer", "expected_snippets"),
    [
        (
            "rich.md",
            _write_long_markdown,
            ["Long Markdown Smoke", "Section 6 body line", "long markdown parser tail"],
        ),
        (
            "rich.html",
            _write_long_html,
            ["Long HTML Smoke", "html parser table note", "Region Delta"],
        ),
        (
            "rich.xlsx",
            _write_wide_xlsx,
            ["Summary", "Region-12", "row-12-note", "North"],
        ),
        (
            "rich.pptx",
            _write_multi_slide_pptx,
            ["PPTX Deck Smoke", "Agenda slide parser smoke", "Closing slide parser smoke"],
        ),
        (
            "rich.pdf",
            _write_complex_pdf,
            ["Complex PDF Smoke", "Second Page Title", "Second page parser smoke body"],
        ),
        (
            "rich.csv",
            lambda path: _write_text(path, "name,score,team\nalpha,7,core\nbeta,9,infra\n"),
            ["name", "alpha", "infra"],
        ),
        (
            "rich.tsv",
            lambda path: _write_text(path, "name\tscore\tteam\nalpha\t7\tcore\nbeta\t9\tinfra\n"),
            ["name\tscore\tteam", "beta\t9\tinfra"],
        ),
        (
            "rich.xml",
            lambda path: _write_text(
                path,
                '<?xml version="1.0"?><report><title>XML Smoke</title><item priority="high">xml parser smoke body</item></report>',
            ),
            ["XML Smoke", "xml parser smoke body"],
        ),
        (
            "rich.rst",
            lambda path: _write_text(
                path,
                "RST Smoke\n=========\n\nrst parser smoke body\n\n* bullet one\n* bullet two\n",
            ),
            ["RST Smoke", "rst parser smoke body", "bullet two"],
        ),
        (
            "rich.eml",
            lambda path: _write_text(
                path,
                'From: parser@example.com\nTo: team@example.com\nSubject: EML Smoke\nMIME-Version: 1.0\nContent-Type: text/plain; charset="utf-8"\n\neml parser smoke body\n',
            ),
            ["Subject: EML Smoke", "eml parser smoke body"],
        ),
    ],
)
def test_doc_parser_regression_for_customer_like_default_path_samples(
    tmp_path: Path,
    filename: str,
    writer: Callable[[Path], None],
    expected_snippets: list[str],
) -> None:
    sample = tmp_path / filename
    writer(sample)

    parser = DocParser(parser_config={"use_markitdown": True, "use_mineru": False})
    parts = parser.parse_file(sample, metadata={"source": filename})
    combined_text = _extract_text(parts)

    assert parts, f"{filename} should produce parsed parts"
    for snippet in expected_snippets:
        assert snippet in combined_text, f"{filename} should contain {snippet!r}, got: {combined_text!r}"
